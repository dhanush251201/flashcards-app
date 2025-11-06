"""
Service for parsing various file formats (PDF, PPT, TXT, DOCX)
Extracts text content from uploaded files for AI processing.
"""

import io
from typing import Optional

from fastapi import HTTPException, UploadFile
from loguru import logger
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document


class FileParserService:
    """Service for parsing various file formats"""

    SUPPORTED_EXTENSIONS = {".pdf", ".ppt", ".pptx", ".txt", ".docx"}
    MAX_FILE_SIZE = 10 * 1024 * 1024  # 10 MB

    @staticmethod
    async def parse_file(file: UploadFile) -> str:
        """
        Parse uploaded file and extract text content.

        Args:
            file: Uploaded file object

        Returns:
            Extracted text content

        Raises:
            HTTPException: If file format is unsupported or parsing fails
        """
        # Validate file extension
        file_ext = FileParserService._get_file_extension(file.filename)
        if file_ext not in FileParserService.SUPPORTED_EXTENSIONS:
            raise HTTPException(
                status_code=400,
                detail=f"Unsupported file format. Supported formats: {', '.join(FileParserService.SUPPORTED_EXTENSIONS)}"
            )

        # Read file content
        try:
            content = await file.read()
            file_size = len(content)

            # Validate file size
            if file_size > FileParserService.MAX_FILE_SIZE:
                raise HTTPException(
                    status_code=400,
                    detail=f"File too large. Maximum size: {FileParserService.MAX_FILE_SIZE / (1024 * 1024)} MB"
                )

            # Parse based on file type
            if file_ext == ".pdf":
                return FileParserService._parse_pdf(content)
            elif file_ext in {".ppt", ".pptx"}:
                return FileParserService._parse_pptx(content)
            elif file_ext == ".txt":
                return FileParserService._parse_txt(content)
            elif file_ext == ".docx":
                return FileParserService._parse_docx(content)
            else:
                raise HTTPException(
                    status_code=400,
                    detail="Unsupported file format"
                )

        except HTTPException:
            raise
        except Exception as e:
            logger.error(f"Error parsing file {file.filename}: {str(e)}")
            raise HTTPException(
                status_code=500,
                detail=f"Failed to parse file: {str(e)}"
            )

    @staticmethod
    def _get_file_extension(filename: Optional[str]) -> str:
        """Extract file extension from filename"""
        if not filename:
            raise HTTPException(status_code=400, detail="Filename is required")

        if "." not in filename:
            raise HTTPException(status_code=400, detail="File must have an extension")

        return "." + filename.rsplit(".", 1)[1].lower()

    @staticmethod
    def _parse_pdf(content: bytes) -> str:
        """Parse PDF file and extract text"""
        try:
            pdf_file = io.BytesIO(content)
            pdf_reader = PdfReader(pdf_file)

            text_content = []
            for page_num, page in enumerate(pdf_reader.pages, 1):
                try:
                    text = page.extract_text()
                    if text.strip():
                        text_content.append(f"--- Page {page_num} ---\n{text}")
                except Exception as e:
                    logger.warning(f"Could not extract text from page {page_num}: {str(e)}")
                    continue

            if not text_content:
                raise ValueError("No text content could be extracted from PDF")

            return "\n\n".join(text_content)

        except Exception as e:
            logger.error(f"Error parsing PDF: {str(e)}")
            raise ValueError(f"Failed to parse PDF: {str(e)}")

    @staticmethod
    def _parse_pptx(content: bytes) -> str:
        """Parse PowerPoint file and extract text"""
        try:
            pptx_file = io.BytesIO(content)
            presentation = Presentation(pptx_file)

            text_content = []
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_text = []

                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)

                if slide_text:
                    text_content.append(
                        f"--- Slide {slide_num} ---\n{chr(10).join(slide_text)}"
                    )

            if not text_content:
                raise ValueError("No text content could be extracted from PowerPoint")

            return "\n\n".join(text_content)

        except Exception as e:
            logger.error(f"Error parsing PowerPoint: {str(e)}")
            raise ValueError(f"Failed to parse PowerPoint: {str(e)}")

    @staticmethod
    def _parse_txt(content: bytes) -> str:
        """Parse text file"""
        try:
            # Try UTF-8 first, fall back to latin-1
            try:
                text = content.decode("utf-8")
            except UnicodeDecodeError:
                text = content.decode("latin-1")

            if not text.strip():
                raise ValueError("Text file is empty")

            return text.strip()

        except Exception as e:
            logger.error(f"Error parsing text file: {str(e)}")
            raise ValueError(f"Failed to parse text file: {str(e)}")

    @staticmethod
    def _parse_docx(content: bytes) -> str:
        """Parse Word document and extract text"""
        try:
            docx_file = io.BytesIO(content)
            document = Document(docx_file)

            text_content = []
            for para_num, paragraph in enumerate(document.paragraphs, 1):
                if paragraph.text.strip():
                    text_content.append(paragraph.text)

            if not text_content:
                raise ValueError("No text content could be extracted from Word document")

            return "\n\n".join(text_content)

        except Exception as e:
            logger.error(f"Error parsing Word document: {str(e)}")
            raise ValueError(f"Failed to parse Word document: {str(e)}")


# Singleton instance
file_parser_service = FileParserService()
